"""Shared fixtures: creates sample test files in a temp directory."""
import io
import struct
import pytest
from pathlib import Path


@pytest.fixture(scope="session")
def sample_dir(tmp_path_factory) -> Path:
    d = tmp_path_factory.mktemp("samples")
    _make_png(d / "sample.png")
    _make_jpg(d / "sample.jpg")
    _make_bmp(d / "sample.bmp")
    _make_gif(d / "sample.gif")
    _make_docx(d / "sample.docx")
    _make_xlsx(d / "sample.xlsx")
    _make_pptx(d / "sample.pptx")
    return d


# --- minimal file builders ---

def _make_png(path: Path):
    from PIL import Image
    img = Image.new("RGB", (100, 100), color=(200, 100, 50))
    img.save(str(path), format="PNG")


def _make_jpg(path: Path):
    from PIL import Image
    img = Image.new("RGB", (100, 100), color=(50, 150, 200))
    img.save(str(path), format="JPEG")


def _make_bmp(path: Path):
    from PIL import Image
    img = Image.new("RGB", (100, 100), color=(100, 200, 100))
    img.save(str(path), format="BMP")


def _make_gif(path: Path):
    from PIL import Image
    img = Image.new("P", (100, 100))
    img.save(str(path), format="GIF")


def _make_docx(path: Path):
    """Minimal valid .docx (ZIP with required OOXML parts)."""
    import zipfile
    content_types = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        '</Types>'
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        '</Relationships>'
    )
    document = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:body><w:p><w:r><w:t>Hello IXP</w:t></w:r></w:p></w:body>'
        '</w:document>'
    )
    word_rels = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
    )
    with zipfile.ZipFile(path, "w") as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", document)
        z.writestr("word/_rels/document.xml.rels", word_rels)


def _make_xlsx(path: Path):
    """Minimal valid .xlsx."""
    import zipfile
    content_types = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/xl/workbook.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
        '<Override PartName="/xl/worksheets/sheet1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>'
        '</Types>'
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="xl/workbook.xml"/>'
        '</Relationships>'
    )
    workbook = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<sheets><sheet name="Sheet1" sheetId="1" r:id="rId1"/></sheets>'
        '</workbook>'
    )
    wb_rels = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" '
        'Target="worksheets/sheet1.xml"/>'
        '</Relationships>'
    )
    sheet = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        '<sheetData><row r="1"><c r="A1" t="inlineStr"><is><t>Hello</t></is></c></row></sheetData>'
        '</worksheet>'
    )
    with zipfile.ZipFile(path, "w") as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("xl/workbook.xml", workbook)
        z.writestr("xl/_rels/workbook.xml.rels", wb_rels)
        z.writestr("xl/worksheets/sheet1.xml", sheet)


def _make_pptx(path: Path):
    """Create a proper PPTX with one slide using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if title:
        title.text = "Hello IXP"
    prs.save(str(path))
